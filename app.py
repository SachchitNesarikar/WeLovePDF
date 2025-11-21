from flask import Flask, render_template, request, send_file
from PyPDF2 import PdfReader, PdfWriter, PdfMerger
from PIL import Image
from io import BytesIO
from werkzeug.utils import secure_filename
import zipfile
import docx2pdf
import tempfile
import os
from dotenv import load_dotenv
import base64
import google.generativeai as genai
from pdfminer.high_level import extract_text
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from docx import Document
from pdf2image import convert_from_path
from flask import Flask, request, jsonify
from pdf2image import convert_from_bytes
from PIL import Image
import pytesseract
import io
import google.genai as genai
from google.genai import types
from pptx import Presentation
from pptx.util import Inches
from generate import get_subtopics, get_notes

load_dotenv()

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
client = genai.Client(api_key=GEMINI_API_KEY)
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = os.path.join(os.getcwd(), '')
pytesseract.pytesseract.tesseract_cmd = r"/usr/bin/tesseract"

def extract_text_with_pdfminer(pdf_bytes):
    with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
        tmp.write(pdf_bytes)
        tmp_path = tmp.name
    return extract_text(tmp_path)

@app.route("/")
def index():
    return render_template("index.html")

@app.route("/image_to_pdf", methods=["POST"])
def image_to_pdf():
    files = request.files.getlist("image_files")
    output_name = request.form.get("output_name", "image_to_pdf").strip()

    image_list = []
    for file in files:
        if file and file.filename:
            try:
                image = Image.open(file.stream).convert("RGB")
                image_list.append(image)
            except Exception as e:
                print(f"Skipping {file.filename}: {e}")

    if not image_list:
        return "No valid images uploaded", 400

    if not output_name.lower().endswith(".pdf"):
        output_name += ".pdf"

    output_stream = BytesIO()
    image_list[0].save(output_stream, save_all=True, append_images=image_list[1:], format="PDF")
    output_stream.seek(0)

    return send_file(output_stream, as_attachment=True, download_name=output_name)

@app.route("/split", methods=["POST"])
def split_pdf():
    file = request.files["pdf_file"]
    start_page = int(request.form.get("start_page")) - 1
    end_page = int(request.form.get("end_page"))
    custom_name = request.form.get("output_name", "split")

    reader = PdfReader(file.stream)
    writer = PdfWriter()

    for i in range(start_page, end_page):
        writer.add_page(reader.pages[i])

    output_stream = BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)

    output_filename = f"{secure_filename(custom_name)}.pdf"
    return send_file(output_stream, as_attachment=True, download_name=output_filename)

@app.route("/merge", methods=["POST"])
def merge_pdfs():
    custom_name = request.form.get("output_name", "merged")
    merger = PdfMerger()
    uploaded_files = request.files.getlist("pdf_files")
    if not uploaded_files or all(f.filename == "" for f in uploaded_files):
        return "No files uploaded", 400
    for file in uploaded_files:
        if file and file.filename:
            merger.append(file.stream)
    output_stream = BytesIO()
    merger.write(output_stream)
    merger.close()
    output_stream.seek(0)
    output_filename = f"{secure_filename(custom_name)}.pdf"
    return send_file(output_stream, as_attachment=True, download_name=output_filename)

@app.route("/compress", methods=["POST"])
def compress_pdf():
    file = request.files["pdf_file"]
    output_name = request.form.get("output_name", "compressed")

    if not file or not file.filename.endswith(".pdf"):
        return "Please upload a valid PDF file", 400

    reader = PdfReader(file.stream)
    writer = PdfWriter()

    for page in reader.pages:
        page.compress_content_streams()
        writer.add_page(page)


    output_stream = BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)

    output_filename = f"{secure_filename(output_name)}.pdf"
    return send_file(output_stream, as_attachment=True, download_name=output_filename)

@app.route("/delete_pages", methods=["POST"])
def delete_pages():
    file = request.files["pdf_file"]
    pages_to_delete = request.form.get("pages_to_delete", "").strip()
    output_name = request.form.get("output_name", "DeletedPages")

    if not file or not file.filename.endswith(".pdf"):
        return "Please upload a valid PDF file", 400
    if not pages_to_delete:
        return "Please specify page numbers to delete", 400

    pages_to_delete_set = set()
    for part in pages_to_delete.split(","):
        part = part.strip()
        if "-" in part:
            start, end = map(int, part.split("-"))
            pages_to_delete_set.update(range(start - 1, end))
        else:
            pages_to_delete_set.add(int(part) - 1)

    reader = PdfReader(file.stream)
    writer = PdfWriter()

    for i, page in enumerate(reader.pages):
        if i not in pages_to_delete_set:
            writer.add_page(page)

    output_stream = BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)

    return send_file(output_stream, as_attachment=True, download_name=f"{secure_filename(output_name)}.pdf")

@app.route("/word_to_pdf", methods=["POST"])
def word_to_pdf():
    word_file = request.files["word_file"]
    output_name = request.form.get("output_name", "ConvertedPDF")

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, secure_filename(word_file.filename))
        output_path = os.path.join(tmpdir, f"{secure_filename(output_name)}.pdf")
        word_file.save(input_path)
        docx2pdf.convert(input_path, output_path)
        with open(output_path, "rb") as f:
            pdf_bytes = f.read()

    return send_file(BytesIO(pdf_bytes), as_attachment=True, download_name=f"{secure_filename(output_name)}.pdf")

@app.route("/pdf_to_word", methods=["POST"])
def pdf_to_word():
    file = request.files["pdf_file"]
    output_name = request.form.get("output_name", "ConvertedWord")

    reader = PdfReader(file.stream)
    text = "\n\n".join([page.extract_text() or "" for page in reader.pages])

    doc_stream = BytesIO()
    doc = Document()
    doc.add_paragraph(text)
    doc.save(doc_stream)
    doc_stream.seek(0)

    return send_file(doc_stream, as_attachment=True, download_name=f"{secure_filename(output_name)}.docx")

@app.route("/compress_image", methods=["POST"])
def compress_image():
    file = request.files["image_file"]
    quality = int(request.form.get("quality", 70))
    output_name = request.form.get("output_name", "CompressedImage")

    image = Image.open(file.stream)
    output_stream = BytesIO()
    image.save(output_stream, format="JPEG", quality=quality)
    output_stream.seek(0)

    return send_file(output_stream, as_attachment=True, download_name=f"{secure_filename(output_name)}.jpg")

@app.route('/ocr_pdf', methods=['POST'])
def ocr_pdf_handler():
    if 'pdf_file' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400
    pdf_bytes = request.files['pdf_file'].read()

    try:
        images = convert_from_bytes(pdf_bytes, dpi=300)
    except Exception as e:
        return jsonify({'error': f'PDF conversion failed: {str(e)}'}), 500

    prompt_parts = [types.Part.from_text("Extract all text with formatting from these PDF pages:")]
    for img in images:
        buffered = BytesIO()
        img.save(buffered, format="PNG")
        img_bytes = buffered.getvalue()
        prompt_parts.append(types.Part.from_bytes(data=img_bytes, mime_type="image/png"))

    try:
        response = client.predict(
            model="gemini-2.5-flash",
            prompt=prompt_parts,
            temperature=0,
            top_k=1,
            top_p=0.95,
            max_output_tokens=1024,
        )
        return jsonify({'text': response.candidates[0].content})
    except Exception as e:
        return jsonify({'error': f'Gemini OCR failed: {str(e)}'}), 500

@app.route("/edit_pdf", methods=["POST"])
def edit_pdf():
    file = request.files.get("pdf_file")
    edit_text = request.form.get("edit_text", "").strip()
    page_number_str = request.form.get("page_number")
    output_name = request.form.get("output_name", "EditedPDF")

    if not file or not file.filename.endswith(".pdf"):
        return "Please upload a valid PDF file", 400
    if not edit_text:
        return "Please enter text to add", 400

    try:
        page_number = int(page_number_str) - 1 if page_number_str else 0

        overlay_stream = BytesIO()
        c = canvas.Canvas(overlay_stream, pagesize=letter)
        c.drawString(72, 720, edit_text)
        c.save()
        overlay_stream.seek(0)

        reader = PdfReader(file.stream)
        overlay_reader = PdfReader(overlay_stream)
        writer = PdfWriter()

        for i, page in enumerate(reader.pages):
            if i == page_number:
                page.merge_page(overlay_reader.pages[0])
            writer.add_page(page)

        output_stream = BytesIO()
        writer.write(output_stream)
        output_stream.seek(0)

        return send_file(output_stream, as_attachment=True, download_name=f"{secure_filename(output_name)}.pdf")
    except Exception as e:
        return f"PDF edit failed: {str(e)}", 500

@app.route("/extract_images", methods=["POST"])
def extract_images():
    pdf_file = request.files.get('pdf_file')
    output_name = request.form.get('output_name') or "ExtractedImages"

    if not pdf_file:
        return "No file uploaded", 400

    pdf_path = os.path.join(app.config['UPLOAD_FOLDER'], pdf_file.filename)
    pdf_file.save(pdf_path)

    with tempfile.TemporaryDirectory() as temp_dir:
        try:
            images = convert_from_path(pdf_path, dpi=200)
        except Exception as e:
            os.remove(pdf_path)
            return f"Error converting PDF: {e}"

        for i, img in enumerate(images):
            img_path = os.path.join(temp_dir, f"page_{i+1}.png")
            img.save(img_path, "PNG")

        zip_path = os.path.join(app.config['UPLOAD_FOLDER'], f"{output_name}.zip")
        with zipfile.ZipFile(zip_path, 'w') as zipf:
            for img_name in os.listdir(temp_dir):
                zipf.write(os.path.join(temp_dir, img_name), img_name)

    os.remove(pdf_path)
    return send_file(zip_path, as_attachment=True)

@app.route("/rotate_pdf", methods=["POST"])
def rotate_pdf():
    file = request.files["pdf_file"]
    rotation_angle = int(request.form.get("rotation_angle", 90))
    output_name = request.form.get("output_name", "RotatedPDF")

    if not file or not file.filename.endswith(".pdf"):
        return "Please upload a valid PDF file", 400

    reader = PdfReader(file.stream)
    writer = PdfWriter()

    for page in reader.pages:
        page.rotate(rotation_angle)
        writer.add_page(page)

    output_stream = BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)

    output_filename = f"{secure_filename(output_name)}.pdf"
    return send_file(output_stream, as_attachment=True, download_name=output_filename)

@app.route("/lock_pdf", methods=["POST"])
def lock_pdf():
    file = request.files.get("pdf_file")
    password = request.form.get("password", "").strip()
    output_name = request.form.get("output_name", "LockedPDF"
                                   )
    if not file or not file.filename.endswith(".pdf"):
        return "Please upload a valid PDF file", 400
    if not password:
        return "Please provide a password", 400
    reader = PdfReader(file.stream)

    writer = PdfWriter()
    for page in reader.pages:
        writer.add_page(page)

    writer.encrypt(password)
    output_stream = BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)

    return send_file(output_stream, as_attachment=True, download_name=f"{secure_filename(output_name)}.pdf")

@app.route("/unlock_pdf", methods=["POST"])
def unlock_pdf():
    file = request.files.get("pdf_file")
    password = request.form.get("password", "").strip()
    output_name = request.form.get("output_name", "UnlockedPDF")

    if not file or not file.filename.endswith(".pdf"):
        return "Please upload a valid PDF file", 400
    if not password:
        return "Please provide the password", 400
    
    reader = PdfReader(file.stream)

    if reader.is_encrypted:
        try:
            reader.decrypt(password)
        except Exception as e:
            return f"Failed to unlock PDF: {str(e)}", 400
        
    writer = PdfWriter()

    for page in reader.pages:
        writer.add_page(page)

    output_stream = BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)

    return send_file(output_stream, as_attachment=True, download_name=f"{secure_filename(output_name)}.pdf")

@app.route("/ppt_to_pdf", methods=["POST"])
def ppt_to_pdf():
    ppt_file = request.files.get("ppt_file")
    output_name = request.form.get("output_name", "ConvertedPPT")

    if not ppt_file or not ppt_file.filename.endswith((".ppt", ".pptx")):
        return "Please upload a valid PowerPoint file", 400

    with tempfile.TemporaryDirectory() as tmpdir:
        input_path = os.path.join(tmpdir, secure_filename(ppt_file.filename))
        output_path = os.path.join(tmpdir, f"{secure_filename(output_name)}.pdf")
        ppt_file.save(input_path)

        try:
            os.system(f'libreoffice --headless --convert-to pdf "{input_path}" --outdir "{tmpdir}"')
        except Exception as e:
            return f"PPT to PDF conversion failed: {str(e)}", 500

        with open(output_path, "rb") as f:
            pdf_bytes = f.read()

    return send_file(BytesIO(pdf_bytes), as_attachment=True, download_name=f"{secure_filename(output_name)}.pdf")

@app.route("/pdf_to_ppt", methods=["POST"])
def pdf_to_ppt():
    pdf_file = request.files.get("pdf_file")
    output_name = request.form.get("output_name", "ConvertedPPT")

    if not pdf_file or not pdf_file.filename.endswith(".pdf"):
        return "Please upload a valid PDF file", 400

    try:
        images = convert_from_bytes(pdf_file.read(), dpi=200)
    except Exception as e:
        return f"PDF to images conversion failed: {str(e)}", 500

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]

    for img in images:
        slide = prs.slides.add_slide(blank_slide_layout)
        img_stream = BytesIO()
        img.save(img_stream, format="PNG")
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)

    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)

    return send_file(ppt_stream, as_attachment=True, download_name=f"{secure_filename(output_name)}.pptx")

@app.route("/crop_pdf", methods=["POST"])
def crop_pdf():
    file = request.files.get("pdf_file")
    crop_values = request.form.get("crop_values", "").strip()
    output_name = request.form.get("output_name", "CroppedPDF")

    if not file or not file.filename.endswith(".pdf"):
        return "Please upload a valid PDF file", 400
    if not crop_values:
        return "Please provide crop values", 400

    try:
        x1, y1, x2, y2 = map(float, crop_values.split(","))
    except:
        return "Invalid crop values format", 400

    reader = PdfReader(file.stream)
    writer = PdfWriter()

    for page in reader.pages:
        page.mediabox.lower_left = (x1, y1)
        page.mediabox.upper_right = (x2, y2)
        writer.add_page(page)

    output_stream = BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)
    return send_file(output_stream, as_attachment=True, download_name=f"{secure_filename(output_name)}.pdf")

@app.route("/sign_pdf", methods=["POST"])
def sign_pdf():
    pdf_file = request.files.get("pdf_file")
    signature_file = request.files.get("signature_file")
    position = request.form.get("signature_position", "100,100,200,80") 
    output_name = request.form.get("output_name", "SignedPDF")

    if not pdf_file or not signature_file:
        return "Please upload PDF and signature image", 400

    try:
        x, y, w, h = map(float, position.split(","))
    except:
        return "Invalid position format", 400

    reader = PdfReader(pdf_file.stream)
    writer = PdfWriter()

    for page in reader.pages:
        packet = BytesIO()
        can = canvas.Canvas(packet, pagesize=letter)
        can.drawImage(signature_file, x, y, width=w, height=h, preserveAspectRatio=True, mask='auto')
        can.save()
        packet.seek(0)

        overlay = PdfReader(packet)
        page.merge_page(overlay.pages[0])
        writer.add_page(page)

    output_stream = BytesIO()
    writer.write(output_stream)
    output_stream.seek(0)

    return send_file(output_stream, as_attachment=True, download_name=f"{secure_filename(output_name)}.pdf")

@app.route("/generate_notes", methods=["POST"])
def generate_notes_route():
    heading = request.form.get("heading").strip()
    
    if not heading:
        return jsonify({"error": "Please provide a heading"}), 400
    
    subtopics = get_subtopics(heading)
    if not subtopics:
        return jsonify({"error": "No subtopics generated. Try a different heading."}), 400
    
    notes = get_notes(subtopics)

    pdf_stream = BytesIO()
    pdf_canvas = canvas.Canvas(pdf_stream, pagesize=letter)
    
    pdf_canvas.setFont("Helvetica-Bold", 16)
    pdf_canvas.drawString(100, 750, f"Topic: {heading}")
    pdf_canvas.setFont("Helvetica", 12)
    
    y_position = 730
    pdf_canvas.drawString(100, y_position, "Subtopics:")
    y_position -= 20
    
    for i, subtopic in enumerate(subtopics, start=1):
        pdf_canvas.drawString(100, y_position, f"{i}. {subtopic}")
        y_position -= 15

    pdf_canvas.showPage()
    pdf_canvas.setFont("Helvetica-Bold", 16)
    pdf_canvas.drawString(100, 750, f"Notes on {heading}")
    pdf_canvas.setFont("Helvetica", 12)
    y_position = 730
    pdf_canvas.drawString(100, y_position, "Notes:")
    y_position -= 20
    
    pdf_canvas.setFont("Helvetica", 10)
    lines = notes.split("\n")
    for line in lines:
        pdf_canvas.drawString(100, y_position, line)
        y_position -= 12
        if y_position < 50:  
            pdf_canvas.showPage()
            y_position = 750  

    pdf_canvas.save()
    pdf_stream.seek(0)

    return send_file(pdf_stream, as_attachment=True, download_name=f"{heading}_notes.pdf")

if __name__ == "__main__":
    app.run(debug=True)
